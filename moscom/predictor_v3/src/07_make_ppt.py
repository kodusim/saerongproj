"""7단계: 발표용 PPT 생성 (심플 디자인).

구성
  1 표지
  2 데이터 소개
  3 사용 데이터 (포집 + 기상)
  4 모델 설계 — 왜 변화량(Δ)을 예측하는가
  5 비교 모델
  6 결과 지표 (R²/RMSE)
  7 예측 기여도
  8 실제값 vs 예측값 (관측소 6개)
  9 규모대별 성능
 10 한계 및 다음 단계
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).parent))
from config import OUT_DIR, ROOT

# --- 디자인 토큰 (심플: 남색 + 회색 2색 체계) ---
NAVY = RGBColor(0x2E, 0x5C, 0x8A)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x77, 0x77, 0x77)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"

W, H = Inches(13.333), Inches(7.5)   # 16:9


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # 빈 레이아웃


def textbox(slide, x, y, w, h, text, size=18, bold=False, color=DARK,
            align=PP_ALIGN.LEFT, spacing=1.0, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = spacing
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = FONT
    return tb


def bar(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def header(slide, title, sub=None):
    """상단 제목 + 남색 언더바."""
    textbox(slide, Inches(0.75), Inches(0.42), Inches(11.8), Inches(0.6),
            title, size=28, bold=True, color=DARK)
    bar(slide, Inches(0.75), Inches(1.12), Inches(0.9), Inches(0.055), NAVY)
    if sub:
        textbox(slide, Inches(0.75), Inches(1.26), Inches(11.8), Inches(0.4),
                sub, size=13, color=GRAY)


def card(slide, x, y, w, h, title, value, note=None):
    """지표 카드 — 제목/값/비고를 세로로 겹치지 않게 배치.

    값은 길이에 따라 폰트를 줄여 한 줄에 들어가게 하고, 비고는 값 박스
    바로 아래에 붙여 카드 높이와 무관하게 겹침이 생기지 않도록 한다.
    """
    bar(slide, x, y, w, h, LIGHT)
    pad = Inches(0.22)
    iw = w - pad * 2
    textbox(slide, x + pad, y + Inches(0.14), iw, Inches(0.28),
            title, size=12, color=GRAY)
    vsize = 30 if len(value) <= 7 else 24 if len(value) <= 9 else 21
    vh = Inches(0.62)
    vy = y + Inches(0.46)
    textbox(slide, x + pad, vy, iw, vh, value, size=vsize, bold=True, color=NAVY)
    if note:
        textbox(slide, x + pad, vy + vh, iw, Inches(0.3), note, size=11, color=GRAY)


def picture_fit(slide, img, x, y, w, h):
    """비율 유지하며 (x,y,w,h) 영역 안에 중앙 배치."""
    from PIL import Image
    iw, ih = Image.open(img).size
    scale = min(w / iw, h / ih)
    nw, nh = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(img, int(x + (w - nw) / 2), int(y + (h - nh) / 2), nw, nh)


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ---------- 1. 표지 ----------
    s = add_slide(prs)
    bar(s, 0, 0, W, Inches(2.55), NAVY)
    textbox(s, Inches(0.95), Inches(0.85), Inches(11), Inches(0.6),
            "모기 발생량 예측 모델", size=40, bold=True, color=WHITE)
    textbox(s, Inches(0.95), Inches(1.62), Inches(11), Inches(0.5),
            "관측소별 야간 포집량 1~3일 앞 예측", size=17, color=RGBColor(0xD5, 0xE2, 0xF0))
    textbox(s, Inches(0.95), Inches(3.15), Inches(11), Inches(1.6),
            "· 원본 측정 로그 127,158건 · 58개 장비 · 13개 권역\n"
            "· 머신러닝 4종 비교 · 기상 데이터 및 장비 상태 결합\n"
            "· 전 관측소의 96% 이상에서 기준선 대비 통계적으로 유의한 개선",
            size=15, color=DARK, spacing=1.75)
    bar(s, Inches(0.95), Inches(5.55), Inches(0.7), Inches(0.045), NAVY)
    textbox(s, Inches(0.95), Inches(5.75), Inches(11), Inches(0.4),
            "2026. 07", size=13, color=GRAY)

    # ---------- 2. 데이터 소개 ----------
    s = add_slide(prs)
    header(s, "데이터 개요", "MOSCOM 포집 장비 원본 측정 로그 (가공 없음) · 야간 수집창 18시~익일 05시 (KST)")
    cw, gap = Inches(2.72), Inches(0.28)
    x0 = Inches(0.75)
    for i, (t, v, n) in enumerate([
            ("원본 측정 로그", "127,158건", "장비별 시간당 1건"),
            ("장비(관측소)", "58대", "13개 권역 (전국)"),
            ("학습 데이터", "5,192건", "장비 × 업무일"),
            ("하루 최대", "20,098마리", "성수기 급증 구간")]):
        card(s, x0 + i * (cw + gap), Inches(1.85), cw, Inches(1.65), t, v, n)

    textbox(s, Inches(0.75), Inches(3.95), Inches(5.6), Inches(0.4),
            "원본 데이터 항목", size=15, bold=True, color=NAVY)
    textbox(s, Inches(0.75), Inches(4.4), Inches(5.6), Inches(2.2),
            "· 측정일시 · 누적 포집 카운트\n"
            "· 카운터 리셋 플래그\n"
            "· 팬 가동 여부 (실제 포집 동작 여부)\n"
            "· 배터리 잔량 (장비 상태 지표)",
            size=13.5, color=DARK, spacing=1.75)

    textbox(s, Inches(6.9), Inches(3.95), Inches(5.7), Inches(0.4),
            "가공본 대비 확보한 정보", size=15, bold=True, color=NAVY)
    textbox(s, Inches(6.9), Inches(4.4), Inches(5.7), Inches(2.2),
            "· 장비 단위 식별 → 동명 관측소 정확히 분리\n"
            "· 미플래그 카운터 롤백 2,846건 추가 검출\n"
            "· 팬 미가동일 310일을 '무발생'이 아닌 '미작동'으로 구분\n"
            "· 측정 누락 시점을 명시적 결측으로 처리",
            size=13.5, color=DARK, spacing=1.75)

    # ---------- 3. 사용 데이터 ----------
    s = add_slide(prs)
    header(s, "투입 데이터", "포집 실측 데이터에 기상 변수를 결합해 학습")

    bar(s, Inches(0.75), Inches(1.9), Inches(5.75), Inches(4.5), LIGHT)
    textbox(s, Inches(1.05), Inches(2.15), Inches(5.2), Inches(0.4),
            "① 모기 포집 원본 로그", size=16, bold=True, color=NAVY)
    textbox(s, Inches(1.05), Inches(2.68), Inches(5.2), Inches(3.4),
            "출처 : MOSCOM 포집 장비 원본 DB\n"
            "기간 : 2026-04-15 ~ 07-28 (105일)\n"
            "범위 : 58대 장비 / 13개 권역\n\n"
            "파생 변수\n"
            "  · 과거 1~7일 포집량 (lag)\n"
            "  · 3일·7일 이동평균 / 변동성\n"
            "  · 평소 대비 편차 (z-score)\n"
            "  · 권역 내 동시 발생 수준\n"
            "  · 팬 가동시간당 포집강도",
            size=13, color=DARK, spacing=1.6)

    bar(s, Inches(6.82), Inches(1.9), Inches(5.75), Inches(4.5), LIGHT)
    textbox(s, Inches(7.12), Inches(2.15), Inches(5.2), Inches(0.4),
            "② 기상 데이터", size=16, bold=True, color=NAVY)
    textbox(s, Inches(7.12), Inches(2.68), Inches(5.2), Inches(3.4),
            "출처 : 기상청 종관기상관측(ASOS) 일자료\n"
            "기간 : 포집 데이터와 동일 구간 전체\n"
            "범위 : 13개 권역 관측지점별 일 단위\n\n"
            "투입 변수\n"
            "  · 기온 (일 평균·최고·최저)\n"
            "  · 야간 수집창 기온·습도·풍속\n"
            "  · 강수량 및 7·14일 누적 강수\n"
            "  · 유효적산온도(GDD) · 무강수 경과일",
            size=13, color=DARK, spacing=1.6)

    textbox(s, Inches(0.75), Inches(6.62), Inches(11.8), Inches(0.4),
            "※ 예측 대상일의 기상은 실제 운영 시 예보값이므로 학습에서 제외 — "
            "기준일까지의 실측 기상만 사용해 정보 누수 차단",
            size=11.5, color=GRAY)

    # ---------- 4. 모델 설계 ----------
    s = add_slide(prs)
    header(s, "모델 설계", "핵심 : 절대량이 아닌 '변화량(Δ)'을 학습")

    textbox(s, Inches(0.75), Inches(1.8), Inches(11.8), Inches(0.4),
            "예측 대상 = log(예측일 포집량) − log(오늘 포집량)   →   '변화량'만 학습",
            size=15, bold=True, color=NAVY)
    bar(s, Inches(0.75), Inches(2.28), Inches(11.8), Inches(0.95), RGBColor(0xE8, 0xF0, 0xF8))
    textbox(s, Inches(1.05), Inches(2.46), Inches(11.2), Inches(0.7),
            "「오늘 값 유지」를 기본값으로 깔고 모델은 그로부터의 편차만 예측.\n"
            "단일 시즌 데이터에서 「시즌이 갈수록 증가」 추세를 외삽해 과대예측하는 문제를 제거.",
            size=13, color=DARK, spacing=1.5)

    # 타깃 변환 방식 비교표 — 실제 실험 결과
    textbox(s, Inches(0.75), Inches(3.42), Inches(11.8), Inches(0.4),
            "타깃 변환 방식 비교 (5개 방식 × 3개 모델 검증)", size=15, bold=True, color=NAVY)
    ty = Inches(3.86)
    bar(s, Inches(0.75), ty, Inches(11.8), Inches(0.42), NAVY)
    for lb, xx, ww in [("변환 방식", 1.0, 4.6), ("로그공간 R²", 6.1, 2.6),
                       ("원공간 R²", 8.9, 2.6)]:
        textbox(s, Inches(xx), ty + Inches(0.06), Inches(ww), Inches(0.33), lb,
                size=12.5, bold=True, color=WHITE)
    rows = [("로그 차분  log(예측일) − log(오늘)", "+0.394", "+0.067", True),
            ("로그 절대량  log(예측일)", "+0.382", "+0.074", False),
            ("원공간 절대량  예측일 마리수", "+0.168", "+0.049", False),
            ("원공간 차분  예측일 − 오늘", "−0.049", "−0.674", False),
            ("배율  예측일 ÷ 오늘", "−0.375", "−3.852", False)]
    for i, (a, b, c, best) in enumerate(rows):
        yy = ty + Inches(0.42) + Inches(0.43) * i
        if best:
            bar(s, Inches(0.75), yy, Inches(11.8), Inches(0.43), RGBColor(0xE8, 0xF0, 0xF8))
        elif i % 2 == 1:
            bar(s, Inches(0.75), yy, Inches(11.8), Inches(0.43), LIGHT)
        col = NAVY if best else DARK
        textbox(s, Inches(1.0), yy + Inches(0.08), Inches(4.9), Inches(0.34),
                ("★ " if best else "     ") + a, size=12, bold=best, color=col)
        textbox(s, Inches(6.1), yy + Inches(0.08), Inches(2.6), Inches(0.34), b,
                size=12, bold=best, color=col)
        textbox(s, Inches(8.9), yy + Inches(0.08), Inches(2.6), Inches(0.34), c,
                size=12, bold=best, color=GRAY if not best else col)

    textbox(s, Inches(0.75), Inches(6.5), Inches(11.8), Inches(0.6),
            "→ 로그 변환이 두 지표 모두에서 우세. 원공간 학습은 극단값(최대 20,098마리)에 끌려가 성능 저하.\n"
            "   추가 적용 : 강한 규제 · 권역 동시성 변수 · 최근 데이터 가중(반감기 21일)",
            size=12, color=GRAY, spacing=1.4)

    # ---------- 5. 비교 모델 ----------
    s = add_slide(prs)
    header(s, "비교 모델", "동일한 검증 체계로 4종 머신러닝 + 기준선 2종 비교")
    rows = [
        ("Ridge 회귀", "선형 모델 + 강한 규제", "1·2일 후 예측 최우수"),
        ("Random Forest", "다수 의사결정나무 평균", "3일 후 예측 최우수"),
        ("XGBoost", "그래디언트 부스팅", "안정적 성능"),
        ("LightGBM", "그래디언트 부스팅", "안정적 성능"),
        ("기준선 : 전일값 유지", "오늘 값을 그대로 사용", "비교 기준 (Baseline)"),
        ("기준선 : 최근 3일 평균", "직전 3일 평균 사용", "비교 기준 (Baseline)"),
    ]
    y = Inches(1.95)
    bar(s, Inches(0.75), y, Inches(11.8), Inches(0.5), NAVY)
    for i, (a, b, c) in enumerate([("모델", "방식", "비고")]):
        textbox(s, Inches(1.0), y + Inches(0.08), Inches(3.4), Inches(0.35), a,
                size=13, bold=True, color=WHITE)
        textbox(s, Inches(4.7), y + Inches(0.08), Inches(4.0), Inches(0.35), b,
                size=13, bold=True, color=WHITE)
        textbox(s, Inches(9.0), y + Inches(0.08), Inches(3.3), Inches(0.35), c,
                size=13, bold=True, color=WHITE)
    for i, (a, b, c) in enumerate(rows):
        yy = y + Inches(0.5) + Inches(0.62) * i
        if i % 2 == 0:
            bar(s, Inches(0.75), yy, Inches(11.8), Inches(0.62), LIGHT)
        is_base = a.startswith("기준선")
        textbox(s, Inches(1.0), yy + Inches(0.14), Inches(3.6), Inches(0.4), a,
                size=13, bold=not is_base, color=GRAY if is_base else DARK)
        textbox(s, Inches(4.7), yy + Inches(0.14), Inches(4.2), Inches(0.4), b,
                size=12.5, color=GRAY)
        textbox(s, Inches(9.0), yy + Inches(0.14), Inches(3.3), Inches(0.4), c,
                size=12.5, color=NAVY if "최우수" in c else GRAY)

    textbox(s, Inches(0.75), Inches(6.35), Inches(11.8), Inches(0.7),
            "검증 방식 : 시간 순서 유지 · 최근 21일 홀드아웃 + 확장창 3-fold 백테스트\n"
            "평가 대상 : 유효 관측 60일 이상 관측소 (학습에는 전체 사용)",
            size=12, color=GRAY, spacing=1.4)

    # ---------- 6. 결과 지표 ----------
    s = add_slide(prs)
    header(s, "결과 지표", "예측 모델 vs 기준선(전일값 유지) · 검증 21일")
    picture_fit(s, str(OUT_DIR / "05_metrics_r2_rmse.png"),
                Inches(0.7), Inches(1.72), Inches(11.9), Inches(3.45))
    cw, gap = Inches(3.72), Inches(0.3)
    for i, (t, v, n) in enumerate([
            ("1일 후 R²", "0.324", "기준선 0.023"),
            ("2일 후 R²", "0.248", "기준선 −0.250"),
            ("3일 후 R²", "0.196", "기준선 −0.275")]):
        card(s, Inches(0.75) + i * (cw + gap), Inches(5.35), cw, Inches(1.42), t, v, n)
    textbox(s, Inches(0.75), Inches(6.95), Inches(11.8), Inches(0.35),
            "※ R²는 규모 편중을 제거한 로그공간 기준 · 관측소 단위 부트스트랩 검정에서 "
            "1·2·3일 후 모두 통계적으로 유의 (p<0.05, 96~98% 관측소에서 개선)", size=11, color=GRAY)

    # ---------- 7. 예측 기여도 ----------
    s = add_slide(prs)
    header(s, "예측 기여도", "어떤 변수가 예측에 실제로 기여했는가 (순열 중요도)")
    picture_fit(s, str(OUT_DIR / "06_feature_importance.png"),
                Inches(0.7), Inches(1.75), Inches(7.6), Inches(5.1))
    textbox(s, Inches(8.55), Inches(2.0), Inches(4.1), Inches(0.4),
            "주요 발견", size=15, bold=True, color=NAVY)
    bar(s, Inches(8.55), Inches(2.5), Inches(4.05), Inches(3.9), LIGHT)
    textbox(s, Inches(8.8), Inches(2.72), Inches(3.6), Inches(3.5),
            "① 평소 대비 상대 위치가 최대 기여\n"
            "   절대 마리수보다 '평소보다\n"
            "   많은가'가 핵심 신호\n\n"
            "② 권역 동시성이 뒤를 이음\n"
            "   같은 날 인근 관측소의 발생 수준\n\n"
            "③ 팬 가동시간당 포집강도\n"
            "   원본 로그에서만 얻을 수 있는\n"
            "   장비 상태 변수가 상위권 진입",
            size=12, color=DARK, spacing=1.5)

    # ---------- 8. 실제 vs 예측 ----------
    s = add_slide(prs)
    header(s, "실제값 vs 예측값", "관측소 6개소 · 1일 후 예측 · 검증 구간 21일")
    picture_fit(s, str(OUT_DIR / "07_station_6.png"),
                Inches(0.55), Inches(1.72), Inches(12.2), Inches(4.75))
    textbox(s, Inches(0.75), Inches(6.62), Inches(11.8), Inches(0.4),
            "발생 수준과 증감 방향을 안정적으로 추종 · 급증 시점의 변화 방향은 포착하나 "
            "정점의 크기는 보수적으로 예측", size=11.5, color=GRAY)

    # ---------- 9. 규모대별 성능 ----------
    s = add_slide(prs)
    header(s, "발생 규모별 성능", "관측소 규모에 따라 효과가 다르게 나타남")
    picture_fit(s, str(OUT_DIR / "08_scale_breakdown.png"),
                Inches(0.7), Inches(1.75), Inches(7.4), Inches(4.9))
    textbox(s, Inches(8.35), Inches(2.0), Inches(4.3), Inches(0.4),
            "해석", size=15, bold=True, color=NAVY)
    bar(s, Inches(8.35), Inches(2.5), Inches(4.25), Inches(3.9), LIGHT)
    textbox(s, Inches(8.6), Inches(2.72), Inches(3.8), Inches(3.5),
            "저발생 구간에서 효과 최대\n"
            "  10마리 이하 구간에서\n"
            "  36~39% 개선\n\n"
            "대형 발생지도 뚜렷한 개선\n"
            "  1000마리 이상에서 28% 개선\n\n"
            "11~50마리 구간은 한계\n"
            "  변동이 불규칙해 기준선이 우세\n"
            "  → 해당 구간 별도 처리 필요",
            size=12, color=DARK, spacing=1.5)

    # ---------- 10. 한계 및 다음 단계 ----------
    s = add_slide(prs)
    header(s, "한계 및 향후 과제")
    textbox(s, Inches(0.75), Inches(1.85), Inches(5.75), Inches(0.4),
            "현재 한계", size=15, bold=True, color=NAVY)
    bar(s, Inches(0.75), Inches(2.35), Inches(5.75), Inches(4.0), LIGHT)
    textbox(s, Inches(1.02), Inches(2.6), Inches(5.2), Inches(3.6),
            "· 단일 시즌(105일) 데이터\n"
            "  전년 동기 정보가 없어 계절 주기 학습 불가\n\n"
            "· 카운터 롤백 구간 과소집계\n"
            "  음수 증분을 0으로 처리 (미플래그 2,846건)\n\n"
            "· 11~50마리 구간 개선 미흡\n\n"
            "· 급증 예측은 방향성 위주\n"
            "  정점의 크기는 보수적으로 예측",
            size=12.5, color=DARK, spacing=1.5)

    textbox(s, Inches(6.82), Inches(1.85), Inches(5.75), Inches(0.4),
            "향후 과제", size=15, bold=True, color=NAVY)
    bar(s, Inches(6.82), Inches(2.35), Inches(5.75), Inches(4.0),
        RGBColor(0xE8, 0xF0, 0xF8))
    textbox(s, Inches(7.09), Inches(2.6), Inches(5.2), Inches(3.6),
            "① 규모별 예측 전략 분리\n"
            "   구간별로 최적 방식을 다르게 적용\n\n"
            "② 급증 탐지 전용 모델 개발\n"
            "   증가 여부를 별도 분류 문제로 학습\n\n"
            "③ 시간대별 데이터 활용\n"
            "   야간 12개 시간대 패턴에서 선행 신호 탐색\n\n"
            "④ 데이터 축적 후 재학습\n"
            "   다음 시즌 확보 시 계절성 반영 가능",
            size=12.5, color=DARK, spacing=1.5)

    out = ROOT / "모기예측_발표자료_v3.pptx"
    prs.save(out)
    print(f"[저장] {out}")
    print(f"       슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장")
    return out


if __name__ == "__main__":
    build()
